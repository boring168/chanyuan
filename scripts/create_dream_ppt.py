from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path("My_Childhood_Dream_vs_My_Dreams_Right_Now_English.pptx")
IMG_DIR = Path("assets/dance")

# Safe fonts for WPS/PowerPoint
FONT_MAIN = "Calibri"
FONT_AUX = "Arial"

# Palette
TEXT = RGBColor(73, 87, 109)  # muted gray-blue
SUB = RGBColor(100, 112, 128)
LINE = RGBColor(224, 215, 231)
CARD_SOFT = RGBColor(255, 252, 249)
LAV = RGBColor(235, 228, 246)
PINK = RGBColor(246, 226, 233)
BEIGE = RGBColor(246, 238, 228)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN_X = 0.9
TITLE_Y = 0.54
TITLE_W = 11.9
CONTENT_Y = 1.72
CONTENT_H = 5.2


def set_run(run, size=20, bold=False, color=TEXT, font=FONT_MAIN):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_full_background(slide, image_name: str, overlay=0.5) -> None:
    """Use dance image only as full-slide background visual."""
    image_path = IMG_DIR / image_name
    if not image_path.exists():
        # Non-empty fallback, still no blank image box.
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(240, 235, 245)
        bg.line.fill.background()
    else:
        slide.shapes.add_picture(str(image_path), Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))

    veil = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
    veil.fill.solid()
    veil.fill.fore_color.rgb = RGBColor(255, 250, 247)
    veil.fill.transparency = overlay
    veil.line.fill.background()

    # cinematic soft gradient band for premium depth
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(1.6))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(255, 252, 249)
    band.fill.transparency = 0.35
    band.line.fill.background()


def add_header(slide, title: str) -> None:
    box = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(TITLE_Y), Inches(TITLE_W), Inches(0.9))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    set_run(p.runs[0], size=34, bold=True)

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(MARGIN_X), Inches(1.33), Inches(2.7), Inches(0.07))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(213, 191, 226)
    bar.line.fill.background()


def add_chapter_tag(slide, idx: str):
    badge = add_content_card(slide, x=11.7, y=0.48, w=0.75, h=0.42, alpha=0.0)
    badge.fill.fore_color.rgb = PINK
    badge.line.fill.background()
    p = badge.text_frame.paragraphs[0]
    p.text = idx
    p.alignment = PP_ALIGN.CENTER
    set_run(p.runs[0], size=12, bold=True, color=SUB, font=FONT_AUX)


def add_content_card(slide, x=MARGIN_X, y=CONTENT_Y, w=11.55, h=CONTENT_H, alpha=0.08):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_SOFT
    card.fill.transparency = alpha
    card.line.color.rgb = LINE
    card.line.width = Pt(1)
    return card


def add_bullets(slide, bullets, x, y, w, h, size=20):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.space_after = Pt(10)
        set_run(p.runs[0], size=size, color=TEXT)


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, "soft_dream_1.jpg", overlay=0.36)
    card = add_content_card(slide, x=1.0, y=1.52, w=11.3, h=4.45, alpha=0.14)
    card.line.fill.background()

    tag = add_content_card(slide, x=4.95, y=1.05, w=3.4, h=0.5, alpha=0.0)
    tag.fill.fore_color.rgb = LAV
    tag.line.fill.background()
    tp = tag.text_frame.paragraphs[0]
    tp.text = "Class Presentation"
    tp.alignment = PP_ALIGN.CENTER
    set_run(tp.runs[0], size=13, color=SUB, font=FONT_AUX)

    title = slide.shapes.add_textbox(Inches(1.35), Inches(2.05), Inches(10.7), Inches(1.85))
    tf_title = title.text_frame
    tf_title.clear()
    p = tf_title.paragraphs[0]
    p.text = "My Childhood Dream"
    p.alignment = PP_ALIGN.CENTER
    set_run(p.runs[0], size=36, bold=True)
    p2 = tf_title.add_paragraph()
    p2.text = "vs. My Dreams Right Now"
    p2.alignment = PP_ALIGN.CENTER
    set_run(p2.runs[0], size=34, bold=True)

    sub = slide.shapes.add_textbox(Inches(2.0), Inches(3.9), Inches(9.3), Inches(1.25))
    tf = sub.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "My dream is still dance, but my understanding has become more mature."
    p1.alignment = PP_ALIGN.CENTER
    set_run(p1.runs[0], size=19, color=SUB)
    p2 = tf.add_paragraph()
    p2.text = "Presented by Meichen Jiang"
    p2.alignment = PP_ALIGN.CENTER
    set_run(p2.runs[0], size=16, color=SUB)


def slide_childhood(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, "soft_dream_2.jpg", overlay=0.62)
    add_header(slide, "My Childhood Dream")
    add_chapter_tag(slide, "01")
    add_content_card(slide)
    add_bullets(
        slide,
        [
            "When I was little, dance felt magical, soft, and full of light.",
            "I loved the stage and imagined myself shining in front of everyone.",
            "My dream was simple and emotional: to dance beautifully and be seen.",
        ],
        x=1.35,
        y=2.25,
        w=10.8,
        h=3.8,
    )


def slide_now(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, "mature_focus_1.jpg", overlay=0.56)
    add_header(slide, "My Dream Right Now")
    add_chapter_tag(slide, "02")
    add_content_card(slide)
    add_bullets(
        slide,
        [
            "As a first-year university student, I still truly love dance.",
            "Now I focus more on practice, discipline, and steady self-improvement.",
            "I want to balance study and passion, and keep dance as part of real growth.",
        ],
        x=1.35,
        y=2.25,
        w=10.8,
        h=3.8,
    )


def slide_similarities(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, "soft_dream_1.jpg", overlay=0.66)
    add_header(slide, "Similarities")
    add_chapter_tag(slide, "03")
    add_content_card(slide, x=0.95, y=1.85, w=11.45, h=5.0)

    items = [
        ("Love for Dance", "Dance has always been close to my heart."),
        ("Self-Expression", "I still express my feelings through movement."),
        ("Confidence & Growth", "I still hope dance helps me become better and more confident."),
    ]
    for i, (head, body) in enumerate(items):
        x = 1.18 + i * 3.77
        card = add_content_card(slide, x=x, y=2.26, w=3.35, h=4.18, alpha=0.04)
        card.line.color.rgb = RGBColor(219, 206, 228)
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 1.5), Inches(2.43), Inches(0.36), Inches(0.36))
        dot.fill.solid()
        dot.fill.fore_color.rgb = LAV
        dot.line.fill.background()
        hb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.62), Inches(2.95), Inches(0.85))
        hp = hb.text_frame.paragraphs[0]
        hp.text = head
        set_run(hp.runs[0], size=20, bold=True)
        bb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(3.55), Inches(2.95), Inches(2.45))
        bp = bb.text_frame.paragraphs[0]
        bp.text = body
        set_run(bp.runs[0], size=17, color=SUB)


def slide_differences(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, "mature_focus_2.jpg", overlay=0.64)
    add_header(slide, "Differences")
    add_chapter_tag(slide, "04")

    left = add_content_card(slide, x=1.2, y=1.96, w=5.42, h=4.86, alpha=0.02)
    left.fill.fore_color.rgb = RGBColor(252, 243, 246)
    left.line.color.rgb = RGBColor(236, 209, 222)
    right = add_content_card(slide, x=6.93, y=1.96, w=5.42, h=4.86, alpha=0.02)
    right.fill.fore_color.rgb = RGBColor(243, 239, 249)
    right.line.color.rgb = RGBColor(213, 204, 236)

    lp = left.text_frame.paragraphs[0]
    lp.text = "Childhood"
    lp.alignment = PP_ALIGN.CENTER
    set_run(lp.runs[0], size=24, bold=True)
    rp = right.text_frame.paragraphs[0]
    rp.text = "Right Now"
    rp.alignment = PP_ALIGN.CENTER
    set_run(rp.runs[0], size=24, bold=True)

    divider = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(6.62), Inches(2.05), Inches(0.22), Inches(4.6))
    divider.fill.solid()
    divider.fill.fore_color.rgb = BEIGE
    divider.fill.transparency = 0.15
    divider.line.fill.background()

    add_bullets(
        slide,
        [
            "More dreamy and emotional",
            "Focused on stage feelings",
            "Driven mostly by imagination",
        ],
        x=1.56,
        y=2.78,
        w=4.75,
        h=3.35,
        size=19,
    )
    add_bullets(
        slide,
        [
            "More realistic and growth-oriented",
            "Focused on consistency and planning",
            "Balancing passion with university life",
        ],
        x=7.28,
        y=2.78,
        w=4.75,
        h=3.35,
        size=19,
    )


def slide_why(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, "mature_focus_1.jpg", overlay=0.6)
    add_header(slide, "Why My Dream Has Changed")
    add_chapter_tag(slide, "05")
    add_content_card(slide)
    add_bullets(
        slide,
        [
            "Growing up made me see dance beyond the stage moment.",
            "More experience taught me that passion needs discipline.",
            "University life made me think more clearly about time and goals.",
            "So the direction did not change, but my mindset became more mature.",
        ],
        x=1.35,
        y=2.2,
        w=10.85,
        h=4.0,
    )


def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, "soft_dream_1.jpg", overlay=0.44)
    add_header(slide, "Final Reflection")
    add_chapter_tag(slide, "06")
    add_content_card(slide, x=1.0, y=1.95, w=11.3, h=4.05, alpha=0.11)

    quote = slide.shapes.add_textbox(Inches(1.58), Inches(2.52), Inches(10.2), Inches(2.5))
    tf = quote.text_frame
    lines = [
        "My dream has always been connected to dance.",
        "What changed is not the direction, but the way I understand it.",
        "Now I want to keep dancing and keep growing at the same time.",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        set_run(p.runs[0], size=24 if i == 1 else 21, bold=(i == 1))

    btn = add_content_card(slide, x=5.02, y=6.12, w=3.3, h=0.7, alpha=0.0)
    btn.fill.fore_color.rgb = LAV
    btn.line.fill.background()
    p = btn.text_frame.paragraphs[0]
    p.text = "Thank you!"
    p.alignment = PP_ALIGN.CENTER
    set_run(p.runs[0], size=21, bold=True)


def build():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide_cover(prs)
    slide_childhood(prs)
    slide_now(prs)
    slide_similarities(prs)
    slide_differences(prs)
    slide_why(prs)
    slide_conclusion(prs)

    prs.save(OUTPUT)


if __name__ == "__main__":
    build()
